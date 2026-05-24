from __future__ import annotations

import subprocess
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path


SUPPORTED_EXTENSIONS = {".md", ".txt", ".text", ".docx", ".pdf", ".ppt", ".pptx"}


# ── .docx (zero-dep: zip + xml) ──────────────────────────────────

def _parse_docx(file_path: Path) -> str:
    """Extract text from .docx using stdlib only (zipfile + xml).

    .docx is a ZIP archive containing word/document.xml with paragraph data.
    """
    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"

    with zipfile.ZipFile(str(file_path), "r") as zf:
        # Main document body
        if "word/document.xml" not in zf.namelist():
            raise ValueError("Invalid .docx: word/document.xml not found")

        tree = ET.parse(zf.open("word/document.xml"))
        root = tree.getroot()

    parts: list[str] = []
    for para in root.iter(f"{ns}p"):
        texts = [node.text for node in para.iter(f"{ns}t") if node.text]
        line = "".join(texts).strip()
        if line:
            parts.append(line)

    return "\n\n".join(parts)


# ── .pptx (zero-dep: zip + xml) ──────────────────────────────────

def _parse_pptx(file_path: Path) -> str:
    """Extract text from .pptx using stdlib only (zipfile + xml).

    .pptx is a ZIP archive; each slide is at ppt/slides/slideN.xml.
    """
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    with zipfile.ZipFile(str(file_path), "r") as zf:
        slide_names = sorted(
            [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
        )

        if not slide_names:
            raise ValueError("Invalid .pptx: no slides found")

        parts: list[str] = []
        for idx, slide_name in enumerate(slide_names, 1):
            tree = ET.parse(zf.open(slide_name))
            root = tree.getroot()

            slide_texts: list[str] = []
            for para in root.iter(f"{ns_a}p"):
                texts = [node.text for node in para.iter(f"{ns_a}t") if node.text]
                line = "".join(texts).strip()
                if line:
                    slide_texts.append(line)

            if slide_texts:
                parts.append(f"[Slide {idx}]\n" + "\n".join(slide_texts))

    return "\n\n".join(parts)


# ── .ppt (legacy binary → textutil fallback) ─────────────────────

def _parse_ppt(file_path: Path) -> str:
    """Extract text from legacy .ppt format.

    Tries macOS textutil first. If unavailable, raises a helpful error.
    """
    # macOS textutil can convert .doc but not .ppt directly.
    # Try python-pptx as optional, otherwise error with guidance.
    try:
        from pptx import Presentation
        prs = Presentation(str(file_path))
        parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
            if slide_texts:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(parts)
    except ImportError:
        pass
    except Exception:
        pass

    raise ValueError(
        "Legacy .ppt format requires conversion. "
        "Please save as .pptx first (open in PowerPoint/WPS → Save As → .pptx), "
        "or install python-pptx: pip install python-pptx"
    )


# ── .pdf (macOS native or pymupdf fallback) ──────────────────────

def _parse_pdf(file_path: Path) -> str:
    """Extract text from .pdf.

    Strategy:
    1. Try pymupdf (fitz) if installed — best quality
    2. Fallback: macOS `osascript` + Quartz filter (zero-dep on macOS)
    3. Fallback: `pdftotext` from poppler if installed
    """
    # Strategy 1: pymupdf
    try:
        import fitz
        doc = fitz.open(str(file_path))
        parts: list[str] = []
        for page in doc:
            text = page.get_text().strip()
            if text:
                parts.append(text)
        doc.close()
        if parts:
            return "\n\n".join(parts)
    except ImportError:
        pass

    # Strategy 2: macOS python3 Quartz (Core Graphics) — zero-dep on macOS
    try:
        result = subprocess.run(
            [
                "python3", "-c",
                "import sys\n"
                "from Quartz import PDFDocument\n"
                "from Foundation import NSURL\n"
                "url = NSURL.fileURLWithPath_(sys.argv[1])\n"
                "doc = PDFDocument.alloc().initWithURL_(url)\n"
                "if doc is None: sys.exit(1)\n"
                "parts = []\n"
                "for i in range(doc.pageCount()):\n"
                "    page = doc.pageAtIndex_(i)\n"
                "    text = page.string()\n"
                "    if text and text.strip(): parts.append(text.strip())\n"
                "print('\\n\\n'.join(parts))\n",
                str(file_path),
            ],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass

    # Strategy 3: pdftotext (poppler)
    try:
        result = subprocess.run(
            ["pdftotext", str(file_path), "-"],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass

    raise ValueError(
        "Could not extract text from PDF. Options:\n"
        "1. Install pymupdf: pip install pymupdf\n"
        "2. Install poppler: brew install poppler (provides pdftotext)\n"
        "3. On macOS, ensure Quartz/pyobjc is available"
    )


# ── Main entry ────────────────────────────────────────────────────

def parse_file(file_path: str) -> str:
    """Read a file and return its text content.

    Supports: .md, .txt, .text, .docx, .pdf, .ppt, .pptx

    .docx and .pptx use Python stdlib only (zipfile + xml).
    .pdf tries pymupdf → macOS Quartz → pdftotext (graceful fallback).
    .ppt (legacy) tries python-pptx if installed, otherwise asks for conversion.

    Raises:
        FileNotFoundError: If file does not exist.
        ValueError: If file extension is not supported or extraction fails.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    suffix = path.suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: {suffix}. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    if suffix == ".docx":
        return _parse_docx(path)
    elif suffix == ".pdf":
        return _parse_pdf(path)
    elif suffix == ".pptx":
        return _parse_pptx(path)
    elif suffix == ".ppt":
        return _parse_ppt(path)

    return path.read_text(encoding="utf-8")


def build_extraction_prompt(content: str) -> dict:
    """Build a prompt for LLM to extract knowledge points from study material.

    Args:
        content: The text content of the study material.

    Returns:
        dict with 'system_prompt' and 'user_prompt' keys.
    """
    system_prompt = (
        "你是一个专业的知识点提取助手。从用户提供的学习资料中提取核心知识点。\n"
        "严格按照要求的 JSON 格式输出，不要输出任何其他内容。"
    )

    user_prompt = f"""请从以下学习资料中提取核心知识点。

## 学习资料内容

{content}

## 提取要求
1. 每个知识点必须包含以下字段：
   - id: 唯一标识符（格式：kp_001, kp_002, ...）
   - title: 知识点名称（简洁明确）
   - definition: 知识点的定义（一句话概括）
   - description: 详细描述（可包含原文中的关键内容）
   - tags: 标签列表（用于分类和检索）
2. 提取所有重要的概念、原理、定义、公式等
3. 每个知识点应该是独立的、可测试的单元
4. description 应尽量保留原文中的关键表述

## 输出格式
输出纯 JSON 数组，每个元素格式如下：
```json
[
  {{
    "id": "kp_001",
    "title": "知识点名称",
    "definition": "一句话定义",
    "description": "详细描述，包含原文关键内容",
    "tags": ["标签1", "标签2"]
  }}
]
```

请直接输出 JSON 数组，不要包含 markdown 代码块标记或其他文字。"""

    return {
        "system_prompt": system_prompt,
        "user_prompt": user_prompt,
    }
