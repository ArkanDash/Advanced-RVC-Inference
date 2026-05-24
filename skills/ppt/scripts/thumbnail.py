#!/usr/bin/env python3
"""
Create thumbnail grids from PowerPoint presentation slides.

Creates a grid layout of slide thumbnails with configurable columns (max 6).
Each grid contains up to cols×(cols+1) images. For presentations with more
slides, multiple numbered grid files are created automatically.

The program outputs the names of all files created.

Output:
- Single grid: {prefix}.jpg (if slides fit in one grid)
- Multiple grids: {prefix}-1.jpg, {prefix}-2.jpg, etc.

Grid limits by column count:
- 3 cols: max 12 slides per grid (3×4)
- 4 cols: max 20 slides per grid (4×5)
- 5 cols: max 30 slides per grid (5×6) [default]
- 6 cols: max 42 slides per grid (6×7)

Usage:
    python thumbnail.py input.pptx [output_prefix] [--cols N] [--outline-placeholders]

Examples:
    python thumbnail.py presentation.pptx
    # Creates: thumbnails.jpg (using default prefix)
    # Outputs:
    #   Created 1 grid(s):
    #     - thumbnails.jpg

    python thumbnail.py large-deck.pptx grid --cols 4
    # Creates: grid-1.jpg, grid-2.jpg, grid-3.jpg
    # Outputs:
    #   Created 3 grid(s):
    #     - grid-1.jpg
    #     - grid-2.jpg
    #     - grid-3.jpg

    python thumbnail.py template.pptx analysis --outline-placeholders
    # Creates thumbnail grids with red outlines around text placeholders
"""

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path

from inventory import extract_text_inventory
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

# Constants
THUMBNAIL_WIDTH = 300  # Fixed thumbnail width in pixels
CONVERSION_DPI = 100  # DPI for PDF to image conversion
MAX_COLS = 6  # Maximum number of columns
DEFAULT_COLS = 5  # Default number of columns
JPEG_QUALITY = 95  # JPEG compression quality

# Grid layout constants
GRID_PADDING = 20  # Padding between thumbnails
BORDER_WIDTH = 2  # Border width around thumbnails
FONT_SIZE_RATIO = 0.12  # Font size as fraction of thumbnail width
LABEL_PADDING_RATIO = 0.4  # Label padding as fraction of font size


def main():
    parser = argparse.ArgumentParser(
        description="Create thumbnail grids from PowerPoint slides."
    )
    parser.add_argument("input", help="Input PowerPoint file (.pptx)")
    parser.add_argument(
        "output_prefix",
        nargs="?",
        default="thumbnails",
        help="Output prefix for image files (default: thumbnails, will create prefix.jpg or prefix-N.jpg)",
    )
    parser.add_argument(
        "--cols",
        type=int,
        default=DEFAULT_COLS,
        help=f"Number of columns (default: {DEFAULT_COLS}, max: {MAX_COLS})",
    )
    parser.add_argument(
        "--outline-placeholders",
        action="store_true",
        help="Outline text placeholders with a colored border",
    )

    args = parser.parse_args()

    cols = min(args.cols, MAX_COLS)
    if args.cols > MAX_COLS:
        print(f"Warning: Columns limited to {MAX_COLS} (requested {args.cols})")

    input_path = Path(args.input)
    if not input_path.is_file() or input_path.suffix.lower() != ".pptx":
        sys.exit(f"Error: Invalid PowerPoint file: {args.input}")

    output_path = Path(f"{args.output_prefix}.jpg")
    print(f"Processing: {args.input}")

    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            placeholder_regions = None
            slide_dimensions = None
            if args.outline_placeholders:
                print("Extracting placeholder regions...")
                placeholder_regions, slide_dimensions = get_placeholder_regions(input_path)
                if placeholder_regions:
                    print(f"Found placeholders on {len(placeholder_regions)} slides")

            prs = Presentation(str(input_path))
            total_slides = len(prs.slides)
            hidden_slides = {
                idx + 1
                for idx, slide in enumerate(prs.slides)
                if slide.element.get("show") == "0"
            }

            hidden_info = f" ({len(hidden_slides)} hidden)" if hidden_slides else ""
            print(f"Found {total_slides} slides{hidden_info}")

            slide_images = convert_to_images(input_path, temp_path, CONVERSION_DPI, total_slides, hidden_slides)
            if not slide_images:
                sys.exit("Error: No slides found")

            grid_files = create_grids(
                slide_images, cols, THUMBNAIL_WIDTH, output_path,
                placeholder_regions, slide_dimensions,
            )

            print(f"Created {len(grid_files)} grid(s):")
            for grid_file in grid_files:
                print(f"  - {grid_file}")

    except RuntimeError as e:
        sys.exit(f"Error: {e}")


def create_hidden_slide_placeholder(size):
    """Create placeholder image for hidden slides."""
    img = Image.new("RGB", size, color="#F0F0F0")
    draw = ImageDraw.Draw(img)
    line_width = max(5, min(size) // 100)
    draw.line([(0, 0), size], fill="#CCCCCC", width=line_width)
    draw.line([(size[0], 0), (0, size[1])], fill="#CCCCCC", width=line_width)
    return img


def get_placeholder_regions(pptx_path):
    """Extract ALL text regions from the presentation.

    Returns a tuple of (placeholder_regions, slide_dimensions).
    text_regions is a dict mapping slide indices to lists of text regions.
    Each region is a dict with 'left', 'top', 'width', 'height' in inches.
    slide_dimensions is a tuple of (width_inches, height_inches).
    """
    prs = Presentation(str(pptx_path))
    inventory = extract_text_inventory(pptx_path, prs)
    placeholder_regions = {}

    slide_width_inches = (prs.slide_width or 9144000) / 914400.0
    slide_height_inches = (prs.slide_height or 5143500) / 914400.0

    for slide_key, shapes in inventory.items():
        slide_idx = int(slide_key.split("-")[1])
        regions = [
            {"left": s.left, "top": s.top, "width": s.width, "height": s.height}
            for s in shapes.values()
        ]
        if regions:
            placeholder_regions[slide_idx] = regions

    return placeholder_regions, (slide_width_inches, slide_height_inches)


def _pptx_to_pdf(pptx_path, temp_dir):
    """Convert PPTX to PDF via LibreOffice. Returns path to the PDF file."""
    pdf_path = temp_dir / f"{pptx_path.stem}.pdf"
    result = subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(temp_dir), str(pptx_path)],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0 or not pdf_path.exists():
        raise RuntimeError("PDF conversion failed")
    return pdf_path


def _pdf_to_images(pdf_path, temp_dir, dpi):
    """Convert PDF pages to JPEG images via pdftoppm. Returns sorted image paths."""
    result = subprocess.run(
        ["pdftoppm", "-jpeg", "-r", str(dpi), str(pdf_path), str(temp_dir / "slide")],
        capture_output=True,
        text=True,
    )
    if result.returncode != 0:
        raise RuntimeError("Image conversion failed")
    return sorted(temp_dir.glob("slide-*.jpg"))


def convert_to_images(pptx_path, temp_dir, dpi, total_slides, hidden_slides):
    """Convert PowerPoint to images via PDF, inserting placeholders for hidden slides."""
    pdf_path = _pptx_to_pdf(pptx_path, temp_dir)
    visible_images = _pdf_to_images(pdf_path, temp_dir, dpi)

    if not visible_images:
        return []

    with Image.open(visible_images[0]) as img:
        placeholder_size = img.size

    all_images = []
    visible_idx = 0
    for slide_num in range(1, total_slides + 1):
        if slide_num in hidden_slides:
            placeholder_path = temp_dir / f"hidden-{slide_num:03d}.jpg"
            create_hidden_slide_placeholder(placeholder_size).save(placeholder_path, "JPEG")
            all_images.append(placeholder_path)
        else:
            if visible_idx < len(visible_images):
                all_images.append(visible_images[visible_idx])
                visible_idx += 1

    return all_images


def create_grids(
    image_paths,
    cols,
    width,
    output_path,
    placeholder_regions=None,
    slide_dimensions=None,
):
    """Create multiple thumbnail grids from slide images, max cols×(cols+1) images per grid."""
    max_images_per_grid = cols * (cols + 1)
    grid_files = []
    total_images = len(image_paths)

    for chunk_idx, start_idx in enumerate(range(0, total_images, max_images_per_grid)):
        chunk_images = image_paths[start_idx: start_idx + max_images_per_grid]

        grid = create_grid(chunk_images, cols, width, start_idx, placeholder_regions, slide_dimensions)

        if total_images <= max_images_per_grid:
            grid_filename = output_path
        else:
            grid_filename = output_path.parent / f"{output_path.stem}-{chunk_idx + 1}{output_path.suffix}"

        grid_filename.parent.mkdir(parents=True, exist_ok=True)
        grid.save(str(grid_filename), quality=JPEG_QUALITY)
        grid_files.append(str(grid_filename))

    return grid_files


def create_grid(
    image_paths,
    cols,
    width,
    start_slide_num=0,
    placeholder_regions=None,
    slide_dimensions=None,
):
    """Create thumbnail grid from slide images with optional placeholder outlining."""
    font_size = int(width * FONT_SIZE_RATIO)
    label_padding = int(font_size * LABEL_PADDING_RATIO)

    with Image.open(image_paths[0]) as img:
        aspect = img.height / img.width
    height = int(width * aspect)

    rows = (len(image_paths) + cols - 1) // cols
    grid_w = cols * width + (cols + 1) * GRID_PADDING
    grid_h = rows * (height + font_size + label_padding * 2) + (rows + 1) * GRID_PADDING

    grid = Image.new("RGB", (grid_w, grid_h), "white")
    draw = ImageDraw.Draw(grid)

    try:
        font = ImageFont.load_default(size=font_size)
    except Exception:
        font = ImageFont.load_default()

    for i, img_path in enumerate(image_paths):
        row, col = i // cols, i % cols
        x = col * width + (col + 1) * GRID_PADDING
        y_base = row * (height + font_size + label_padding * 2) + (row + 1) * GRID_PADDING

        label = f"{start_slide_num + i}"
        bbox = draw.textbbox((0, 0), label, font=font)
        text_w = bbox[2] - bbox[0]
        draw.text((x + (width - text_w) // 2, y_base + label_padding), label, fill="black", font=font)

        y_thumbnail = y_base + label_padding + font_size + label_padding

        with Image.open(img_path) as img:
            orig_w, orig_h = img.size

            if placeholder_regions and (start_slide_num + i) in placeholder_regions:
                if img.mode != "RGBA":
                    img = img.convert("RGBA")

                regions = placeholder_regions[start_slide_num + i]
                if slide_dimensions:
                    slide_w_in, slide_h_in = slide_dimensions
                else:
                    slide_w_in = orig_w / CONVERSION_DPI
                    slide_h_in = orig_h / CONVERSION_DPI

                x_scale = orig_w / slide_w_in
                y_scale = orig_h / slide_h_in

                overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
                overlay_draw = ImageDraw.Draw(overlay)
                stroke_width = max(5, min(orig_w, orig_h) // 150)

                for region in regions:
                    px_left = int(region["left"] * x_scale)
                    px_top = int(region["top"] * y_scale)
                    px_right = px_left + int(region["width"] * x_scale)
                    px_bottom = px_top + int(region["height"] * y_scale)
                    overlay_draw.rectangle(
                        [(px_left, px_top), (px_right, px_bottom)],
                        outline=(255, 0, 0, 255),
                        width=stroke_width,
                    )

                img = Image.alpha_composite(img, overlay).convert("RGB")

            img.thumbnail((width, height), Image.Resampling.LANCZOS)
            w, h = img.size
            tx = x + (width - w) // 2
            ty = y_thumbnail + (height - h) // 2
            grid.paste(img, (tx, ty))

            if BORDER_WIDTH > 0:
                draw.rectangle(
                    [(tx - BORDER_WIDTH, ty - BORDER_WIDTH), (tx + w + BORDER_WIDTH - 1, ty + h + BORDER_WIDTH - 1)],
                    outline="gray",
                    width=BORDER_WIDTH,
                )

    return grid


if __name__ == "__main__":
    main()
