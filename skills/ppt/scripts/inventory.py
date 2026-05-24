#!/usr/bin/env python3
"""
Extract structured text content from PowerPoint presentations.

Usage:
    python inventory.py input.pptx output.json [--issues-only]
"""

import argparse
import json
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.shapes.base import BaseShape

# Public type alias used by replace.py: slide_id -> {shape_id -> ShapeData}
InventoryData = Dict[str, Dict[str, "ShapeData"]]

_EMU = 914400  # EMUs per inch
_BULLET_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_ALIGN_MAP = {
    PP_ALIGN.CENTER: "CENTER",
    PP_ALIGN.RIGHT: "RIGHT",
    PP_ALIGN.JUSTIFY: "JUSTIFY",
}


def _is_cjk(ch: str) -> bool:
    """True for full-width CJK characters (Chinese, Japanese, Korean, full-width forms)."""
    cp = ord(ch)
    return (
        0x4E00 <= cp <= 0x9FFF   # CJK Unified Ideographs
        or 0x3400 <= cp <= 0x4DBF  # CJK Extension A
        or 0x3040 <= cp <= 0x30FF  # Hiragana / Katakana
        or 0xFF00 <= cp <= 0xFFEF  # Full-width ASCII & half-width Katakana
        or 0xAC00 <= cp <= 0xD7AF  # Hangul syllables
    )


class ParagraphData:
    """Text and formatting for one paragraph."""

    def __init__(self, paragraph: Any):
        self.text: str = paragraph.text.strip()
        self.bullet: bool = False
        self.level: Optional[int] = None
        self.alignment: Optional[str] = None
        self.space_before: Optional[float] = None
        self.space_after: Optional[float] = None
        self.font_name: Optional[str] = None
        self.font_size: Optional[float] = None
        self.bold: Optional[bool] = None
        self.italic: Optional[bool] = None
        self.underline: Optional[bool] = None
        self.color: Optional[str] = None
        self.theme_color: Optional[str] = None
        self.line_spacing: Optional[float] = None

        # Bullet detection
        pPr = getattr(getattr(paragraph, "_p", None), "pPr", None)
        if pPr is not None and (
            pPr.find(f"{_BULLET_NS}buChar") is not None
            or pPr.find(f"{_BULLET_NS}buAutoNum") is not None
        ):
            self.bullet = True
            self.level = getattr(paragraph, "level", None)

        # Alignment (omit LEFT — it's the default)
        align = getattr(paragraph, "alignment", None)
        if align in _ALIGN_MAP:
            self.alignment = _ALIGN_MAP[align]

        # Spacing
        sb = getattr(paragraph, "space_before", None)
        if sb:
            self.space_before = sb.pt
        sa = getattr(paragraph, "space_after", None)
        if sa:
            self.space_after = sa.pt

        # Font from first run
        if paragraph.runs:
            font = paragraph.runs[0].font
            self.font_name = font.name or None
            self.font_size = font.size.pt if font.size else None
            self.bold = font.bold
            self.italic = font.italic
            self.underline = font.underline
            try:
                self.color = str(font.color.rgb) if font.color.rgb else None
            except (AttributeError, TypeError):
                try:
                    tc = font.color.theme_color
                    self.theme_color = tc.name if tc else None
                except (AttributeError, TypeError):
                    pass

        # Line spacing (after font so font_size is available)
        ls = getattr(paragraph, "line_spacing", None)
        if ls is not None:
            if hasattr(ls, "pt"):
                self.line_spacing = round(ls.pt, 2)
            else:
                # Multiplier — convert to points using current font size
                self.line_spacing = round(ls * (self.font_size or 12.0), 2)

    def to_dict(self) -> Dict[str, Any]:
        d: Dict[str, Any] = {"text": self.text}
        if self.bullet:
            d["bullet"] = True
            if self.level is not None:
                d["level"] = self.level
        if self.alignment:
            d["alignment"] = self.alignment
        for key in ("space_before", "space_after", "font_size", "line_spacing"):
            val = getattr(self, key)
            if val is not None:
                d[key] = val
        if self.font_name:
            d["font_name"] = self.font_name
        for key in ("bold", "italic", "underline"):
            val = getattr(self, key)
            if val is not None:
                d[key] = val
        if self.color:
            d["color"] = self.color
        elif self.theme_color:
            d["theme_color"] = self.theme_color
        return d


class ShapeData:
    """Position, formatting metadata, and text content for one shape."""

    def __init__(
        self,
        shape: BaseShape,
        absolute_left: Optional[int] = None,
        absolute_top: Optional[int] = None,
        slide: Optional[Any] = None,
    ):
        self.shape = shape
        self.shape_id: str = ""  # assigned after sorting

        # Slide dimensions (for overflow checking)
        self.slide_width_emu: Optional[int] = None
        self.slide_height_emu: Optional[int] = None
        if slide:
            try:
                prs_xml = slide.part.package.presentation_part.presentation
                self.slide_width_emu = prs_xml.slide_width
                self.slide_height_emu = prs_xml.slide_height
            except (AttributeError, TypeError):
                pass

        # Placeholder metadata
        self.placeholder_type: Optional[str] = None
        self.default_font_size: Optional[float] = None
        if getattr(shape, "is_placeholder", False):
            pf = shape.placeholder_format  # type: ignore
            if pf and pf.type:
                self.placeholder_type = str(pf.type).split(".")[-1].split(" ")[0]
                if slide and hasattr(slide, "slide_layout"):
                    self.default_font_size = _layout_font_size(shape, slide.slide_layout)

        # Position in inches (use absolute coords for shapes inside groups)
        left_emu = absolute_left if absolute_left is not None else getattr(shape, "left", 0)
        top_emu = absolute_top if absolute_top is not None else getattr(shape, "top", 0)
        self.left = round(left_emu / _EMU, 2)
        self.top = round(top_emu / _EMU, 2)
        self.width = round(getattr(shape, "width", 0) / _EMU, 2)
        self.height = round(getattr(shape, "height", 0) / _EMU, 2)

        # EMU positions kept for overflow arithmetic
        self.left_emu = left_emu
        self.top_emu = top_emu
        self.width_emu = getattr(shape, "width", 0)
        self.height_emu = getattr(shape, "height", 0)

        # Issue detection
        self.frame_overflow_bottom: Optional[float] = None
        self.slide_overflow_right: Optional[float] = None
        self.slide_overflow_bottom: Optional[float] = None
        self.overlapping_shapes: Dict[str, float] = {}
        self.warnings: List[str] = []
        self._estimate_frame_overflow()
        self._calculate_slide_overflow()
        self._detect_bullet_issues()

    # ------------------------------------------------------------------
    # Issue detection helpers
    # ------------------------------------------------------------------

    def _default_font_size_pts(self) -> float:
        """Best-effort default font size from theme styles."""
        if self.default_font_size:
            return self.default_font_size
        try:
            master = self.shape.part.slide_layout.slide_master  # type: ignore
            style = "titleStyle" if (self.placeholder_type and "TITLE" in self.placeholder_type) else "bodyStyle"
            for child in master.element.iter():
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag == style:
                    for elem in child.iter():
                        if "sz" in elem.attrib:
                            return int(elem.attrib["sz"]) / 100.0
        except Exception:
            pass
        return 14.0  # conservative fallback

    def _estimate_frame_overflow(self) -> None:
        """Estimate text overflow via character-count heuristic (no external deps)."""
        if not hasattr(self.shape, "text_frame"):
            return
        tf = self.shape.text_frame  # type: ignore
        if not tf or not tf.paragraphs:
            return

        # Usable area after text frame margins
        def e2i(v: Any) -> float:
            return (v or 0) / _EMU

        margin_h = e2i(tf.margin_top) + e2i(tf.margin_bottom)
        margin_w = e2i(tf.margin_left) + e2i(tf.margin_right)
        if margin_h == 0:
            margin_h = 0.10  # PowerPoint default: ~0.05" top + 0.05" bottom
        if margin_w == 0:
            margin_w = 0.20  # PowerPoint default: ~0.1" left + 0.1" right
        usable_w = self.width - margin_w
        usable_h = self.height - margin_h
        if usable_w <= 0 or usable_h <= 0:
            return

        default_size = self._default_font_size_pts()
        total_h = 0.0

        for para in tf.paragraphs:
            if not para.text.strip():
                continue
            pd = ParagraphData(para)
            size_pt = pd.font_size or default_size

            # Estimate text width: CJK chars ≈ 1.0× font_size pts, others ≈ 0.5×
            text_w_pts = sum(
                size_pt if _is_cjk(c) else size_pt * 0.5
                for c in para.text
            )
            usable_w_pts = usable_w * 72.0
            n_lines = max(1, -(-int(text_w_pts) // max(1, int(usable_w_pts))))  # ceiling div

            line_h_in = (pd.line_spacing or size_pt) / 72.0
            total_h += (pd.space_before or 0) / 72.0
            total_h += n_lines * line_h_in
            total_h += (pd.space_after or 0) / 72.0

        if total_h > usable_h + 0.05:  # ignore sub-0.05" rounding noise
            self.frame_overflow_bottom = round(total_h - usable_h, 2)

    def _calculate_slide_overflow(self) -> None:
        if self.slide_width_emu is None or self.slide_height_emu is None:
            return
        r = self.left_emu + self.width_emu - self.slide_width_emu
        if r > 0:
            v = round(r / _EMU, 2)
            if v > 0.01:
                self.slide_overflow_right = v
        b = self.top_emu + self.height_emu - self.slide_height_emu
        if b > 0:
            v = round(b / _EMU, 2)
            if v > 0.01:
                self.slide_overflow_bottom = v

    def _detect_bullet_issues(self) -> None:
        if not hasattr(self.shape, "text_frame"):
            return
        for para in self.shape.text_frame.paragraphs:  # type: ignore
            text = para.text.strip()
            if text and any(text.startswith(s + " ") for s in ("•", "●", "○")):
                self.warnings.append("manual_bullet_symbol: use proper bullet formatting")
                break

    # ------------------------------------------------------------------
    # Public interface
    # ------------------------------------------------------------------

    @property
    def paragraphs(self) -> List[ParagraphData]:
        if not hasattr(self.shape, "text_frame"):
            return []
        return [ParagraphData(p) for p in self.shape.text_frame.paragraphs if p.text.strip()]  # type: ignore

    @property
    def has_any_issues(self) -> bool:
        return bool(
            self.frame_overflow_bottom is not None
            or self.slide_overflow_right is not None
            or self.slide_overflow_bottom is not None
            or self.overlapping_shapes
            or self.warnings
        )

    def to_dict(self) -> Dict[str, Any]:
        d: Dict[str, Any] = {
            "left": self.left, "top": self.top,
            "width": self.width, "height": self.height,
        }
        if self.placeholder_type:
            d["placeholder_type"] = self.placeholder_type
        if self.default_font_size:
            d["default_font_size"] = self.default_font_size

        overflow: Dict[str, Any] = {}
        if self.frame_overflow_bottom is not None:
            overflow["frame"] = {"overflow_bottom": self.frame_overflow_bottom}
        slide_ov: Dict[str, float] = {}
        if self.slide_overflow_right is not None:
            slide_ov["overflow_right"] = self.slide_overflow_right
        if self.slide_overflow_bottom is not None:
            slide_ov["overflow_bottom"] = self.slide_overflow_bottom
        if slide_ov:
            overflow["slide"] = slide_ov
        if overflow:
            d["overflow"] = overflow
        if self.overlapping_shapes:
            d["overlap"] = {"overlapping_shapes": self.overlapping_shapes}
        if self.warnings:
            d["warnings"] = self.warnings
        d["paragraphs"] = [p.to_dict() for p in self.paragraphs]
        return d


# ------------------------------------------------------------------
# Module-level helpers
# ------------------------------------------------------------------

def _layout_font_size(shape: BaseShape, slide_layout: Any) -> Optional[float]:
    """Extract default font size from the matching layout placeholder."""
    try:
        shape_type = shape.placeholder_format.type  # type: ignore
        for ph in slide_layout.placeholders:
            if ph.placeholder_format.type == shape_type:
                for elem in ph.element.iter():
                    if "defRPr" in elem.tag and (sz := elem.get("sz")):
                        return float(sz) / 100.0
                break
    except Exception:
        pass
    return None


def _is_valid_shape(shape: BaseShape) -> bool:
    """True if shape has meaningful text and is not a slide-number placeholder."""
    if not hasattr(shape, "text_frame"):
        return False
    tf = shape.text_frame  # type: ignore
    if not tf or not tf.text.strip():
        return False
    if getattr(shape, "is_placeholder", False):
        pf = shape.placeholder_format  # type: ignore
        if pf and pf.type:
            pt = str(pf.type).split(".")[-1].split(" ")[0]
            if pt == "SLIDE_NUMBER":
                return False
            if pt == "FOOTER" and tf.text.strip().isdigit():
                return False
    return True


def _collect_shapes(shape: BaseShape, parent_left: int = 0, parent_top: int = 0):
    """Yield (shape, abs_left, abs_top) tuples, recursing into GroupShapes."""
    if hasattr(shape, "shapes"):  # GroupShape
        g_left = parent_left + getattr(shape, "left", 0)
        g_top = parent_top + getattr(shape, "top", 0)
        for child in shape.shapes:  # type: ignore
            yield from _collect_shapes(child, g_left, g_top)
    elif _is_valid_shape(shape):
        yield (
            shape,
            parent_left + getattr(shape, "left", 0),
            parent_top + getattr(shape, "top", 0),
        )


def _sort_by_position(shapes: List[ShapeData]) -> List[ShapeData]:
    """Sort shapes top-to-bottom, left-to-right (0.5" row tolerance)."""
    if not shapes:
        return shapes
    shapes = sorted(shapes, key=lambda s: (s.top, s.left))
    result: List[ShapeData] = []
    row = [shapes[0]]
    row_top = shapes[0].top
    for s in shapes[1:]:
        if abs(s.top - row_top) <= 0.5:
            row.append(s)
        else:
            result.extend(sorted(row, key=lambda s: s.left))
            row = [s]
            row_top = s.top
    result.extend(sorted(row, key=lambda s: s.left))
    return result


def _detect_overlaps(shapes: List[ShapeData]) -> None:
    """Populate overlapping_shapes for all pairs with meaningful overlap."""
    for i, s1 in enumerate(shapes):
        for s2 in shapes[i + 1:]:
            ow = min(s1.left + s1.width, s2.left + s2.width) - max(s1.left, s2.left)
            oh = min(s1.top + s1.height, s2.top + s2.height) - max(s1.top, s2.top)
            if ow > 0.05 and oh > 0.05:
                area = round(ow * oh, 2)
                s1.overlapping_shapes[s2.shape_id] = area
                s2.overlapping_shapes[s1.shape_id] = area


# ------------------------------------------------------------------
# Public API
# ------------------------------------------------------------------

def extract_text_inventory(
    pptx_path: Path,
    prs: Optional[Any] = None,
    issues_only: bool = False,
) -> InventoryData:
    """Extract text from all slides.

    Returns {slide-N: {shape-N: ShapeData}}, shapes sorted by visual position.
    Pass an existing Presentation object via `prs` to avoid re-loading.
    """
    if prs is None:
        prs = Presentation(str(pptx_path))

    inventory: InventoryData = {}

    for slide_idx, slide in enumerate(prs.slides):
        raw = list(_collect_shapes_from_slide(slide))
        if not raw:
            continue

        shape_data_list = [ShapeData(s, al, at, slide) for s, al, at in raw]
        sorted_shapes = _sort_by_position(shape_data_list)

        for idx, sd in enumerate(sorted_shapes):
            sd.shape_id = f"shape-{idx}"

        if len(sorted_shapes) > 1:
            _detect_overlaps(sorted_shapes)

        if issues_only:
            sorted_shapes = [sd for sd in sorted_shapes if sd.has_any_issues]
        if not sorted_shapes:
            continue

        inventory[f"slide-{slide_idx}"] = {sd.shape_id: sd for sd in sorted_shapes}

    return inventory


def _collect_shapes_from_slide(slide):
    """Yield (shape, abs_left, abs_top) for all valid text shapes on a slide."""
    for shape in slide.shapes:  # type: ignore
        yield from _collect_shapes(shape)


def save_inventory(inventory: InventoryData, output_path: Path) -> None:
    """Serialize inventory to a JSON file."""
    json_data = {
        slide_key: {k: sd.to_dict() for k, sd in shapes.items()}
        for slide_key, shapes in inventory.items()
    }
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(json_data, f, indent=2, ensure_ascii=False)


def main() -> None:
    parser = argparse.ArgumentParser(description="Extract text inventory from a PowerPoint file.")
    parser.add_argument("input", help="Input .pptx file")
    parser.add_argument("output", help="Output .json file")
    parser.add_argument("--issues-only", action="store_true",
                        help="Include only shapes with overflow/overlap issues")
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: File not found: {args.input}")
        sys.exit(1)
    if input_path.suffix.lower() != ".pptx":
        print("Error: Input must be a .pptx file")
        sys.exit(1)

    try:
        inventory = extract_text_inventory(input_path, issues_only=args.issues_only)
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        save_inventory(inventory, output_path)

        total = sum(len(v) for v in inventory.values())
        if args.issues_only:
            print(f"Found {total} shapes with issues across {len(inventory)} slides → {args.output}")
        else:
            print(f"Found {total} text shapes across {len(inventory)} slides → {args.output}")
    except Exception as e:
        import traceback
        print(f"Error: {e}")
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
